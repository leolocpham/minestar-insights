# =============================================================================
# presentation_builder.py – Build a PowerPoint report from MineStar analysis
# =============================================================================
from __future__ import annotations
import io
from datetime import datetime
from typing import List, Dict, Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
NAVY    = RGBColor(0x1a, 0x3a, 0x5c)
COPPER  = RGBColor(0xb8, 0x73, 0x33)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LGRAY   = RGBColor(0xf0, 0xf4, 0xf8)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
GREEN   = RGBColor(0x2e, 0x7d, 0x32)
RED     = RGBColor(0xc6, 0x28, 0x28)
AMBER   = RGBColor(0xe6, 0x5c, 0x00)

SEV_COLOR = {"critical": RED, "warning": AMBER, "info": NAVY, "positive": GREEN}
SEV_HEX   = {"critical": "#c62828", "warning": "#e65c00", "info": "#1a3a5c", "positive": "#2e7d32"}
SEV_LABEL = {"critical": "CRITICAL", "warning": "WARNING", "info": "INFO", "positive": "POSITIVE"}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.85)
ACCENT_H = Inches(0.05)
MARGIN   = Inches(0.3)
CONTENT_TOP = HEADER_H + ACCENT_H + Inches(0.15)
CONTENT_H   = SLIDE_H - CONTENT_TOP - MARGIN

MPLC_NAVY   = "#1a3a5c"
MPLC_COPPER = "#b87333"
MPLC_GREEN  = "#2e7d32"
MPLC_RED    = "#c62828"
MPLC_AMBER  = "#e65c00"
MPLC_LGRAY  = "#f0f4f8"


# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill_rgb: RGBColor, border=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not border:
        shape.line.fill.background()
    return shape


def _textbox(slide, text, left, top, width, height,
             size=12, bold=False, color=DGRAY, align=PP_ALIGN.LEFT,
             word_wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _header(slide, title: str, mine_name: str):
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, NAVY)
    _rect(slide, 0, HEADER_H, SLIDE_W, ACCENT_H, COPPER)
    _textbox(slide, title, MARGIN, Inches(0.12), Inches(9.5), Inches(0.65),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _textbox(slide, mine_name, Inches(9.9), Inches(0.20), Inches(3.2), Inches(0.45),
             size=11, bold=False, color=COPPER, align=PP_ALIGN.RIGHT)


def _chart_img(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _add_image(slide, buf: io.BytesIO, left, top, width, height):
    slide.shapes.add_picture(buf, left, top, width, height)


def _findings_box(slide, items: list, left, top, width, height):
    """Render a list of {severity, text} bullets in a box."""
    _rect(slide, left, top, width, height, LGRAY)
    y = top + Inches(0.15)
    line_h = Inches(0.38)
    for item in items[:8]:
        color = SEV_COLOR.get(item.get("severity", "info"), NAVY)
        tag = SEV_LABEL.get(item.get("severity", "info"), "")
        # Severity tag
        _rect(slide, left + Inches(0.1), y, Inches(0.9), Inches(0.28), color)
        _textbox(slide, tag, left + Inches(0.1), y, Inches(0.9), Inches(0.28),
                 size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Finding text
        _textbox(slide, item.get("text", ""), left + Inches(1.1), y,
                 width - Inches(1.3), Inches(0.3),
                 size=9, bold=False, color=DGRAY)
        y += line_h
        if y + line_h > top + height:
            break


# ---------------------------------------------------------------------------
# Chart builders (matplotlib → BytesIO PNG)
# ---------------------------------------------------------------------------

def _chart_cycle_by_machine(bm: pd.DataFrame) -> io.BytesIO:
    df = bm.head(12).sort_values("avg_cycle")
    fig, ax = plt.subplots(figsize=(7, 4))
    fig.patch.set_facecolor("white")
    ax.set_facecolor(MPLC_LGRAY)
    colors = [MPLC_RED if v > df["avg_cycle"].median() * 1.2 else MPLC_NAVY
              for v in df["avg_cycle"]]
    ax.barh(df["Machine"].astype(str), df["avg_cycle"], color=colors, edgecolor="none")
    ax.set_xlabel("Avg Cycle Time (min)", fontsize=9)
    ax.set_title("Average Cycle Time by Machine", fontsize=11, fontweight="bold",
                 color=MPLC_NAVY, pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8)
    ax.grid(axis="x", alpha=0.3, linestyle="--")
    fig.tight_layout()
    return _chart_img(fig)


def _chart_time_breakdown(bd: dict) -> io.BytesIO:
    labels_map = {
        "queue_time": "Queue", "load_time": "Load",
        "travel_loaded": "Travel (Loaded)", "travel_empty": "Travel (Empty)",
        "dump_time": "Dump", "spot_time": "Spot",
    }
    labels = [labels_map[k] for k in bd if bd[k] and bd[k] > 0]
    values = [bd[k] for k in bd if bd[k] and bd[k] > 0]
    colors = [MPLC_RED if "Queue" in l else MPLC_NAVY if "Travel" in l
              else MPLC_COPPER for l in labels]

    fig, ax = plt.subplots(figsize=(5, 3.5))
    fig.patch.set_facecolor("white")
    wedges, texts, autotexts = ax.pie(
        values, labels=labels, autopct="%1.0f%%", colors=colors,
        startangle=140, pctdistance=0.8,
        textprops={"fontsize": 8},
    )
    for at in autotexts:
        at.set_fontsize(7)
        at.set_color("white")
    ax.set_title("Cycle Time Breakdown", fontsize=10, fontweight="bold",
                 color=MPLC_NAVY, pad=6)
    fig.tight_layout()
    return _chart_img(fig)


def _chart_payload_dist(distribution, target) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(7, 3.8))
    fig.patch.set_facecolor("white")
    ax.set_facecolor(MPLC_LGRAY)
    ax.hist(distribution, bins=30, color=MPLC_NAVY, edgecolor="white",
            linewidth=0.4, alpha=0.85)
    if target:
        ax.axvline(target, color=MPLC_GREEN, linewidth=2, linestyle="--",
                   label=f"Target: {target:.0f} t")
        ax.axvline(target * 0.90, color=MPLC_AMBER, linewidth=1.5, linestyle=":",
                   label="−10% threshold")
        ax.axvline(target * 1.10, color=MPLC_RED, linewidth=1.5, linestyle=":",
                   label="+10% threshold")
        ax.legend(fontsize=8)
    ax.set_xlabel("Payload (t)", fontsize=9)
    ax.set_ylabel("Frequency", fontsize=9)
    ax.set_title("Payload Distribution", fontsize=11, fontweight="bold",
                 color=MPLC_NAVY, pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8)
    ax.grid(axis="y", alpha=0.3, linestyle="--")
    fig.tight_layout()
    return _chart_img(fig)


def _chart_payload_by_machine(bm: pd.DataFrame, target: Optional[float]) -> io.BytesIO:
    df = bm.head(12).sort_values("avg_payload")
    fig, ax = plt.subplots(figsize=(7, 4))
    fig.patch.set_facecolor("white")
    ax.set_facecolor(MPLC_LGRAY)
    colors = [MPLC_GREEN if (target and v >= target * 0.95) else MPLC_AMBER
              for v in df["avg_payload"]]
    ax.barh(df["Machine"].astype(str), df["avg_payload"], color=colors, edgecolor="none")
    if target:
        ax.axvline(target, color=MPLC_NAVY, linewidth=2, linestyle="--",
                   label=f"Target {target:.0f} t")
        ax.legend(fontsize=8)
    ax.set_xlabel("Avg Payload (t)", fontsize=9)
    ax.set_title("Average Payload by Machine", fontsize=11, fontweight="bold",
                 color=MPLC_NAVY, pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8)
    ax.grid(axis="x", alpha=0.3, linestyle="--")
    fig.tight_layout()
    return _chart_img(fig)


def _chart_utilization(bm: pd.DataFrame) -> io.BytesIO:
    has_av = "Availability %" in bm.columns
    has_ut = "Utilization %"  in bm.columns
    has_stacked = all(c in bm.columns for c in ["Operating Hrs", "Idle Hrs", "Down Hrs"])

    if has_stacked:
        df = bm.sort_values("Operating Hrs", ascending=False).head(12)
        fig, ax = plt.subplots(figsize=(7, 4))
        fig.patch.set_facecolor("white")
        ax.set_facecolor(MPLC_LGRAY)
        m = df["Machine"].astype(str)
        ax.barh(m, df["Operating Hrs"], color=MPLC_GREEN, label="Operating", edgecolor="none")
        idle = df.get("Idle Hrs", pd.Series(np.zeros(len(df))))
        ax.barh(m, idle, left=df["Operating Hrs"], color=MPLC_AMBER, label="Idle", edgecolor="none")
        dn = df.get("Down Hrs", pd.Series(np.zeros(len(df))))
        ax.barh(m, dn, left=df["Operating Hrs"] + idle, color=MPLC_RED, label="Down", edgecolor="none")
        ax.legend(fontsize=8, loc="lower right")
        ax.set_xlabel("Hours", fontsize=9)
        ax.set_title("Fleet Hours by Machine", fontsize=11, fontweight="bold",
                     color=MPLC_NAVY, pad=8)
    elif has_av or has_ut:
        col = "Availability %" if has_av else "Utilization %"
        df = bm.sort_values(col).head(14)
        fig, ax = plt.subplots(figsize=(7, 4))
        fig.patch.set_facecolor("white")
        ax.set_facecolor(MPLC_LGRAY)
        colors = [MPLC_RED if v < 75 else MPLC_AMBER if v < 85 else MPLC_GREEN
                  for v in df[col]]
        ax.barh(df["Machine"].astype(str), df[col], color=colors, edgecolor="none")
        ax.axvline(85 if has_av else 70, color=MPLC_NAVY, linewidth=2,
                   linestyle="--", label="Target")
        ax.set_xlim(0, 110)
        ax.legend(fontsize=8)
        ax.set_xlabel(col, fontsize=9)
        ax.set_title(f"{col} by Machine", fontsize=11, fontweight="bold",
                     color=MPLC_NAVY, pad=8)
    else:
        fig, ax = plt.subplots(figsize=(5, 3))
        ax.text(0.5, 0.5, "Insufficient data", ha="center", va="center",
                transform=ax.transAxes, fontsize=12, color="gray")

    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8)
    ax.grid(axis="x", alpha=0.3, linestyle="--")
    fig.tight_layout()
    return _chart_img(fig)


def _chart_operators(by_op: pd.DataFrame) -> io.BytesIO:
    df = by_op.head(12).sort_values("Cycles")
    fig, ax = plt.subplots(figsize=(7, 4))
    fig.patch.set_facecolor("white")
    ax.set_facecolor(MPLC_LGRAY)
    bar_colors = [MPLC_COPPER if i < 3 else MPLC_NAVY
                  for i in range(len(df) - 1, -1, -1)][::-1]
    ax.barh(df["Operator"].astype(str), df["Cycles"],
            color=bar_colors, edgecolor="none")
    ax.set_xlabel("Total Cycles", fontsize=9)
    ax.set_title("Cycles Completed by Operator", fontsize=11, fontweight="bold",
                 color=MPLC_NAVY, pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8)
    ax.grid(axis="x", alpha=0.3, linestyle="--")
    fig.tight_layout()
    return _chart_img(fig)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _slide_title(prs, mine_name: str, date_str: str):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(slide, 0, Inches(3.6), SLIDE_W, Inches(0.06), COPPER)
    _textbox(slide, "MineStar Operations Analysis",
             Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, mine_name,
             Inches(1.2), Inches(3.1), Inches(11), Inches(0.7),
             size=24, bold=False, color=COPPER, align=PP_ALIGN.CENTER)
    _textbox(slide, f"Report generated: {date_str}",
             Inches(1.2), Inches(3.9), Inches(11), Inches(0.5),
             size=12, bold=False, color=RGBColor(0xaa, 0xbb, 0xcc),
             align=PP_ALIGN.CENTER)
    _textbox(slide, "Powered by MineStar Insights App",
             Inches(1.2), Inches(6.8), Inches(11), Inches(0.4),
             size=9, bold=False, color=RGBColor(0x77, 0x88, 0x99),
             align=PP_ALIGN.CENTER)


def _slide_summary(prs, analysis: dict, insights: list, mine_name: str):
    slide = _blank_slide(prs)
    _header(slide, "Executive Summary", mine_name)

    ct = analysis.get("cycle_times", {})
    pl = analysis.get("payload", {})
    ut = analysis.get("utilization", {})
    op = analysis.get("operators", {})

    kpis = [
        ("Total Records",   f"{analysis.get('row_count', 0):,}", ""),
        ("Avg Cycle Time",  f"{ct['avg_cycle']:.1f}" if ct.get("available") else "N/A", "min"),
        ("Avg Payload",     f"{pl['avg_payload']:.1f}" if pl.get("available") else "N/A", "t"),
        ("Fleet Avail.",    f"{ut['avg_availability']:.0f}" if ut.get("avg_availability") else "N/A", "%"),
        ("Operators",       f"{op['operator_count']}" if op.get("available") else "N/A", ""),
        ("Insights Found",  f"{len(insights)}", ""),
    ]

    card_w = Inches(1.9)
    card_h = Inches(1.3)
    gap    = Inches(0.25)
    start_x = Inches(0.35)
    card_top = CONTENT_TOP + Inches(0.1)

    for i, (label, value, unit) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        _rect(slide, x, card_top, card_w, card_h, LGRAY)
        _rect(slide, x, card_top, card_w, Inches(0.07), COPPER)
        _textbox(slide, f"{value} {unit}".strip(),
                 x, card_top + Inches(0.12), card_w, Inches(0.65),
                 size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _textbox(slide, label,
                 x, card_top + Inches(0.78), card_w, Inches(0.4),
                 size=10, bold=False, color=DGRAY, align=PP_ALIGN.CENTER)

    # Top findings list
    top_findings = insights[:6]
    findings_top = card_top + card_h + Inches(0.3)
    findings_h   = SLIDE_H - findings_top - MARGIN

    _rect(slide, MARGIN, findings_top, Inches(12.7), findings_h, LGRAY)
    _textbox(slide, "Key Findings",
             MARGIN + Inches(0.15), findings_top + Inches(0.1),
             Inches(4), Inches(0.35),
             size=12, bold=True, color=NAVY)

    y = findings_top + Inches(0.5)
    for ins in top_findings:
        color = SEV_COLOR.get(ins["severity"], NAVY)
        bullet_text = f"[{SEV_LABEL.get(ins['severity'], '')}]  {ins['finding']}"
        tb = slide.shapes.add_textbox(
            MARGIN + Inches(0.2), y, Inches(12.2), Inches(0.32))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"■  "
        run.font.size = Pt(9)
        run.font.color.rgb = color
        run2 = p.add_run()
        run2.text = ins["finding"]
        run2.font.size = Pt(9)
        run2.font.color.rgb = DGRAY
        y += Inches(0.34)


def _slide_cycle_times(prs, ct: dict, insights: list, mine_name: str):
    slide = _blank_slide(prs)
    _header(slide, "Fleet Cycle Times", mine_name)

    chart_w = Inches(7.8)
    chart_h = Inches(4.2)
    chart_l = MARGIN
    chart_t = CONTENT_TOP + Inches(0.1)

    bm = ct.get("by_machine")
    bd = ct.get("time_breakdown", {})

    if bm is not None and len(bm) > 0:
        buf = _chart_cycle_by_machine(bm)
        _add_image(slide, buf, chart_l, chart_t, chart_w, chart_h)
    elif bd:
        buf = _chart_time_breakdown(bd)
        _add_image(slide, buf, chart_l, chart_t, chart_w, chart_h)

    # KPI strip below chart
    kpi_top = chart_t + chart_h + Inches(0.12)
    kpis = [
        ("Avg Cycle", f"{ct['avg_cycle']:.1f} min"),
        ("Median",    f"{ct['median_cycle']:.1f} min"),
        ("Total Cycles", f"{ct['total_cycles']:,}"),
        ("Best",      f"{ct['min_cycle']:.1f} min"),
    ]
    kw = Inches(1.7)
    kx = chart_l
    for label, val in kpis:
        _rect(slide, kx, kpi_top, kw, Inches(0.65), RGBColor(0xe8, 0xee, 0xf4))
        _textbox(slide, val, kx, kpi_top + Inches(0.02), kw, Inches(0.38),
                 size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _textbox(slide, label, kx, kpi_top + Inches(0.38), kw, Inches(0.25),
                 size=8, color=DGRAY, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.1)

    # Findings panel (right side)
    panel_l = chart_l + chart_w + Inches(0.2)
    panel_w = SLIDE_W - panel_l - MARGIN
    panel_h = CONTENT_H

    _rect(slide, panel_l, CONTENT_TOP, panel_w, panel_h, LGRAY)
    _textbox(slide, "Findings & Actions",
             panel_l + Inches(0.1), CONTENT_TOP + Inches(0.1), panel_w - Inches(0.2),
             Inches(0.35), size=11, bold=True, color=NAVY)

    ct_insights = [i for i in insights if i["category"] == "Cycle Times"]
    y = CONTENT_TOP + Inches(0.5)
    for ins in ct_insights[:5]:
        color = SEV_COLOR.get(ins["severity"], NAVY)
        _rect(slide, panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22), color)
        _textbox(slide, SEV_LABEL[ins["severity"]],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22),
                 size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.24)
        _textbox(slide, ins["finding"],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.55),
                 size=8, color=DGRAY)
        y += Inches(0.57)
        _textbox(slide, f"Action: {ins['recommendation']}",
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.65),
                 size=7, color=COPPER)
        y += Inches(0.7)
        if y > CONTENT_TOP + panel_h - Inches(0.4):
            break


def _slide_payload(prs, pl: dict, insights: list, mine_name: str):
    slide = _blank_slide(prs)
    _header(slide, "Payload Performance", mine_name)

    chart_w = Inches(7.8)
    chart_h = Inches(3.8)
    chart_l = MARGIN
    chart_t = CONTENT_TOP + Inches(0.1)

    bm  = pl.get("by_machine")
    tgt = pl.get("target_payload")
    dist = pl.get("distribution")

    if bm is not None and len(bm) > 0:
        buf = _chart_payload_by_machine(bm, tgt)
        _add_image(slide, buf, chart_l, chart_t, chart_w, chart_h)
    elif dist is not None and len(dist) > 0:
        buf = _chart_payload_dist(dist, tgt)
        _add_image(slide, buf, chart_l, chart_t, chart_w, chart_h)

    # KPI strip
    kpi_top = chart_t + chart_h + Inches(0.12)
    kpis = [
        ("Avg Payload",    f"{pl['avg_payload']:.1f} t"),
        ("Total Tonnes",   f"{pl['total_tonnes']:,.0f} t"),
        ("Total Loads",    f"{pl['total_loads']:,}"),
    ]
    if tgt:
        kpis.append(("Target",     f"{tgt:.1f} t"))
    if pl.get("pct_overloaded") is not None:
        kpis.append(("Overloaded",  f"{pl['pct_overloaded']:.0f}%"))
    if pl.get("pct_underloaded") is not None:
        kpis.append(("Underloaded", f"{pl['pct_underloaded']:.0f}%"))

    kw = Inches(1.9)
    kx = chart_l
    for label, val in kpis[:5]:
        _rect(slide, kx, kpi_top, kw, Inches(0.65), RGBColor(0xe8, 0xee, 0xf4))
        _textbox(slide, val, kx, kpi_top + Inches(0.02), kw, Inches(0.38),
                 size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _textbox(slide, label, kx, kpi_top + Inches(0.38), kw, Inches(0.25),
                 size=8, color=DGRAY, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.08)

    # Findings panel
    panel_l = chart_l + chart_w + Inches(0.2)
    panel_w = SLIDE_W - panel_l - MARGIN
    panel_h = CONTENT_H

    _rect(slide, panel_l, CONTENT_TOP, panel_w, panel_h, LGRAY)
    _textbox(slide, "Findings & Actions",
             panel_l + Inches(0.1), CONTENT_TOP + Inches(0.1), panel_w - Inches(0.2),
             Inches(0.35), size=11, bold=True, color=NAVY)

    pl_insights = [i for i in insights if i["category"] == "Payload"]
    y = CONTENT_TOP + Inches(0.5)
    for ins in pl_insights[:4]:
        color = SEV_COLOR.get(ins["severity"], NAVY)
        _rect(slide, panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22), color)
        _textbox(slide, SEV_LABEL[ins["severity"]],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22),
                 size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.24)
        _textbox(slide, ins["finding"],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.55),
                 size=8, color=DGRAY)
        y += Inches(0.57)
        _textbox(slide, f"Action: {ins['recommendation']}",
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.65),
                 size=7, color=COPPER)
        y += Inches(0.72)
        if y > CONTENT_TOP + panel_h - Inches(0.4):
            break


def _slide_utilization(prs, ut: dict, insights: list, mine_name: str):
    slide = _blank_slide(prs)
    _header(slide, "Equipment Utilization", mine_name)

    chart_w = Inches(7.8)
    chart_h = Inches(4.2)
    chart_l = MARGIN
    chart_t = CONTENT_TOP + Inches(0.1)

    bm = ut.get("by_machine")
    if bm is not None and len(bm) > 0:
        buf = _chart_utilization(bm)
        _add_image(slide, buf, chart_l, chart_t, chart_w, chart_h)

    kpi_top = chart_t + chart_h + Inches(0.12)
    kpis = []
    if ut.get("avg_availability"):
        kpis.append(("Avg Availability", f"{ut['avg_availability']:.1f}%"))
    if ut.get("avg_utilization"):
        kpis.append(("Avg Utilization",  f"{ut['avg_utilization']:.1f}%"))
    if ut.get("total_operating_hours"):
        kpis.append(("Total Op. Hours",  f"{ut['total_operating_hours']:,.0f} h"))
    kw = Inches(2.2)
    kx = chart_l
    for label, val in kpis[:4]:
        _rect(slide, kx, kpi_top, kw, Inches(0.65), RGBColor(0xe8, 0xee, 0xf4))
        _textbox(slide, val, kx, kpi_top + Inches(0.02), kw, Inches(0.38),
                 size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _textbox(slide, label, kx, kpi_top + Inches(0.38), kw, Inches(0.25),
                 size=8, color=DGRAY, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.12)

    panel_l = chart_l + chart_w + Inches(0.2)
    panel_w = SLIDE_W - panel_l - MARGIN
    panel_h = CONTENT_H

    _rect(slide, panel_l, CONTENT_TOP, panel_w, panel_h, LGRAY)
    _textbox(slide, "Findings & Actions",
             panel_l + Inches(0.1), CONTENT_TOP + Inches(0.1), panel_w - Inches(0.2),
             Inches(0.35), size=11, bold=True, color=NAVY)

    ut_insights = [i for i in insights if i["category"] == "Utilization"]
    y = CONTENT_TOP + Inches(0.5)
    for ins in ut_insights[:4]:
        color = SEV_COLOR.get(ins["severity"], NAVY)
        _rect(slide, panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22), color)
        _textbox(slide, SEV_LABEL[ins["severity"]],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22),
                 size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.24)
        _textbox(slide, ins["finding"],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.55),
                 size=8, color=DGRAY)
        y += Inches(0.57)
        _textbox(slide, f"Action: {ins['recommendation']}",
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.65),
                 size=7, color=COPPER)
        y += Inches(0.72)
        if y > CONTENT_TOP + panel_h - Inches(0.4):
            break


def _slide_operators(prs, op: dict, insights: list, mine_name: str):
    slide = _blank_slide(prs)
    _header(slide, "Operator Performance", mine_name)

    chart_w = Inches(7.8)
    chart_h = Inches(4.2)
    chart_l = MARGIN
    chart_t = CONTENT_TOP + Inches(0.1)

    by_op = op.get("by_operator")
    if by_op is not None and len(by_op) > 0:
        buf = _chart_operators(by_op)
        _add_image(slide, buf, chart_l, chart_t, chart_w, chart_h)

    kpi_top = chart_t + chart_h + Inches(0.12)
    kpis = [
        ("Total Operators", str(op.get("operator_count", "N/A"))),
        ("Top Performer",   str(op.get("top_performers", ["N/A"])[0])),
    ]
    if by_op is not None and "Avg Payload (t)" in by_op.columns:
        kpis.append(("Best Avg Payload",
                     f"{by_op['Avg Payload (t)'].max():.1f} t"))
    kw = Inches(2.2)
    kx = chart_l
    for label, val in kpis[:4]:
        _rect(slide, kx, kpi_top, kw, Inches(0.65), RGBColor(0xe8, 0xee, 0xf4))
        _textbox(slide, val, kx, kpi_top + Inches(0.02), kw, Inches(0.38),
                 size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _textbox(slide, label, kx, kpi_top + Inches(0.38), kw, Inches(0.25),
                 size=8, color=DGRAY, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.12)

    panel_l = chart_l + chart_w + Inches(0.2)
    panel_w = SLIDE_W - panel_l - MARGIN
    panel_h = CONTENT_H

    _rect(slide, panel_l, CONTENT_TOP, panel_w, panel_h, LGRAY)
    _textbox(slide, "Findings & Actions",
             panel_l + Inches(0.1), CONTENT_TOP + Inches(0.1), panel_w - Inches(0.2),
             Inches(0.35), size=11, bold=True, color=NAVY)

    op_insights = [i for i in insights if i["category"] == "Operators"]
    y = CONTENT_TOP + Inches(0.5)
    for ins in op_insights[:4]:
        color = SEV_COLOR.get(ins["severity"], NAVY)
        _rect(slide, panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22), color)
        _textbox(slide, SEV_LABEL[ins["severity"]],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.22),
                 size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.24)
        _textbox(slide, ins["finding"],
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.55),
                 size=8, color=DGRAY)
        y += Inches(0.57)
        _textbox(slide, f"Action: {ins['recommendation']}",
                 panel_l + Inches(0.1), y, panel_w - Inches(0.2), Inches(0.65),
                 size=7, color=COPPER)
        y += Inches(0.72)
        if y > CONTENT_TOP + panel_h - Inches(0.4):
            break


def _slide_recommendations(prs, insights: list, mine_name: str):
    slide = _blank_slide(prs)
    _header(slide, "Key Recommendations & Action Plan", mine_name)

    # Sort: critical first
    order = {"critical": 0, "warning": 1, "info": 2, "positive": 3}
    sorted_ins = sorted(insights, key=lambda x: order.get(x["severity"], 9))

    y = CONTENT_TOP + Inches(0.1)
    card_h = Inches(0.95)
    card_gap = Inches(0.1)
    two_col = len(sorted_ins) > 5

    for i, ins in enumerate(sorted_ins[:10]):
        col = 0 if (not two_col or i % 2 == 0) else 1
        row = i if not two_col else i // 2
        card_w = Inches(6.1) if two_col else Inches(12.7)
        x = MARGIN if col == 0 else MARGIN + card_w + Inches(0.4)
        y_card = CONTENT_TOP + Inches(0.1) + row * (card_h + card_gap)

        if y_card + card_h > SLIDE_H - MARGIN:
            break

        sev_color = SEV_COLOR.get(ins["severity"], NAVY)
        _rect(slide, x, y_card, card_w, card_h, LGRAY)
        _rect(slide, x, y_card, Inches(0.07), card_h, sev_color)

        # Category + severity badge
        _textbox(slide, f"{ins['category'].upper()}  —  {SEV_LABEL[ins['severity']]}",
                 x + Inches(0.15), y_card + Inches(0.06),
                 card_w - Inches(0.25), Inches(0.25),
                 size=8, bold=True, color=sev_color)

        # Finding summary
        _textbox(slide, ins["finding"],
                 x + Inches(0.15), y_card + Inches(0.3),
                 card_w - Inches(0.25), Inches(0.3),
                 size=8, color=DGRAY)

        # Action
        _textbox(slide, f"Action: {ins['recommendation'][:160]}",
                 x + Inches(0.15), y_card + Inches(0.6),
                 card_w - Inches(0.25), Inches(0.3),
                 size=7, color=COPPER)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def build_presentation(
    analysis: dict,
    insights: list,
    mine_name: str = "Pinto Valley Mine",
) -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    date_str = datetime.now().strftime("%B %d, %Y")

    _slide_title(prs, mine_name, date_str)
    _slide_summary(prs, analysis, insights, mine_name)

    ct = analysis.get("cycle_times", {})
    pl = analysis.get("payload", {})
    ut = analysis.get("utilization", {})
    op = analysis.get("operators", {})

    if ct.get("available"):
        _slide_cycle_times(prs, ct, insights, mine_name)
    if pl.get("available"):
        _slide_payload(prs, pl, insights, mine_name)
    if ut.get("available"):
        _slide_utilization(prs, ut, insights, mine_name)
    if op.get("available"):
        _slide_operators(prs, op, insights, mine_name)

    _slide_recommendations(prs, insights, mine_name)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
